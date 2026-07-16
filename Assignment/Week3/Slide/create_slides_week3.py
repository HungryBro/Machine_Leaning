"""สร้างสไลด์ภาษาไทยสำหรับ Bias–Variance Lab.

รันจากโฟลเดอร์ใดก็ได้:
    python3 Assignment/Week3/Slide/create_slides_week3.py

จุดสำคัญคือภาพทุกภาพถูกวางแบบ contain โดยคงสัดส่วนเดิม
จึงไม่ถูกยืดหรือถูกตัดขอบเมื่อใส่ลงในสไลด์ 16:9
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PROJECT_DIR = BASE_DIR.parent
PLOTS_DIR = PROJECT_DIR / "plots"
OUTPUT = BASE_DIR / "slides_Week3.pptx"

SLIDE_W, SLIDE_H = 13.333, 7.5


def rgb(hex_value):
    hex_value = hex_value.lstrip("#")
    return RGBColor.from_string(hex_value.upper())


NAVY = rgb("2D6A4F")
INK = rgb("212529")
MUTED = rgb("6C757D")
TEAL = rgb("2D6A4F")
ORANGE = rgb("F77F00")
GREEN = rgb("2D6A4F")
RED = rgb("D62828")
LIGHT = rgb("E9ECEF")
PALE_TEAL = rgb("F8F9FA")
PALE_ORANGE = rgb("F8F9FA")
BG = rgb("F8F9FA")
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_font(run, size, color=INK, bold=False, name="Arial"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        set_font(run, size, color, bold)
    return box


def add_rect(slide, x, y, w, h, fill, radius=False, line=None, line_width=Pt(1)):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_width
    return shape


def add_footer(slide, page):
    add_rect(slide, 0.62, 6.98, 12.08, 0.012, rgb("DEE2E6"))
    add_text(slide, "Machine Learning — Week 3  |  Bias-Variance Lab", 0.62, 7.08, 6.8, 0.2,
             size=8, color=MUTED, margin=0)
    add_text(slide, str(page), 12.12, 7.06, 0.58, 0.22,
             size=9, color=MUTED, bold=True, align=PP_ALIGN.RIGHT, margin=0)


def add_header(slide, title, page, kicker=None):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_rect(slide, 0, 0, SLIDE_W, 0.42, GREEN)
    if kicker:
        add_text(slide, kicker, 0.62, 0.08, 5.8, 0.2,
                 size=10, color=WHITE, bold=True, margin=0)
    else:
        add_text(slide, "Machine Learning — Week 3", 0.62, 0.08, 5.8, 0.2,
                 size=10, color=WHITE, bold=True, margin=0)
    add_text(slide, title, 0.62, 0.58, 12.0, 0.52,
             size=25, color=INK, bold=True, margin=0)
    add_rect(slide, 0.62, 1.08, 0.86, 0.035, RED)


def add_card(slide, x, y, w, h, title, body, accent=TEAL,
             fill=LIGHT, title_size=14, body_size=16, radius=False):
    add_rect(slide, x, y, w, h, fill, radius=radius)
    add_rect(slide, x, y, 0.08, h, accent, radius=False)
    add_text(slide, title, x + 0.22, y + 0.16, w - 0.4, 0.28,
             size=title_size, color=accent, bold=True, margin=0)
    add_text(slide, body, x + 0.22, y + 0.52, w - 0.4, h - 0.62,
             size=body_size, color=INK, margin=0)


def add_bullets(slide, items, x, y, w, h, size=18, color=INK, gap=10):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.space_after = Pt(gap)
        for run in p.runs:
            set_font(run, size, color)
    return box


def add_contained_image(slide, image_path, x, y, w, h):
    """วางภาพในกรอบโดยรักษา aspect ratio และจัดกึ่งกลางเสมอ."""
    with Image.open(image_path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        draw_w = w
        draw_h = w / image_ratio
    else:
        draw_h = h
        draw_w = h * image_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(draw_x), Inches(draw_y),
                             width=Inches(draw_w), height=Inches(draw_h))


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GREEN)
    add_rect(slide, 0.78, 1.52, 1.15, 0.035, WHITE)
    add_text(slide, "BIAS–VARIANCE LAB", 0.78, 1.82, 5.0, 0.28,
             size=12, color=WHITE, bold=True, margin=0)
    add_text(slide, "Bias² และ Variance\nผ่านภาพและการจำลอง", 0.78, 2.22, 11.2, 1.2,
             size=32, color=WHITE, bold=True, margin=0)
    add_text(slide, "ทำไมโมเดลที่ซับซ้อนกว่า ไม่ได้แปลว่าดีกว่าเสมอ", 0.8, 3.78, 10.6, 0.38,
             size=18, color=WHITE, margin=0)
    add_text(slide, "sin(πx)  •  x²  •  Learning Curves", 0.8, 5.55, 7.0, 0.3,
             size=13, color=WHITE, margin=0)
    add_text(slide, "รายงานผลการทดลอง  |  Machine Learning — Week 3", 0.8, 6.1, 8.5, 0.25,
             size=10, color=WHITE, margin=0)
    add_rect(slide, 8.15, 5.05, 0.025, 1.48, WHITE)
    add_text(slide, "ผู้จัดทำ", 8.45, 5.08, 3.8, 0.25,
             size=11, color=WHITE, bold=True, margin=0)
    add_text(slide, "นายสุภมงคล ชอบรัมย์\nนายชินวัตร กิตต๊ะ\nนายกฤตชัย พรายศรี", 8.45, 5.42, 4.15, 0.9,
             size=10, color=WHITE, margin=0)
    return slide


def add_setup_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "โจทย์และการตั้งค่าการทดลอง", 2, "01  Setup")
    add_text(slide, "เราต้องการตอบคำถามเดียว: เมื่อข้อมูลมีจำกัด\nโมเดลแบบไหน generalize ได้ดีที่สุด?",
             0.72, 1.45, 5.45, 1.2, size=23, color=NAVY, bold=True, margin=0)
    add_text(slide, "ออกแบบการทดลองให้เห็น trade-off ระหว่าง bias และ variance",
             0.72, 2.86, 5.1, 0.45, size=15, color=MUTED, margin=0)
    add_card(slide, 6.55, 1.35, 2.85, 1.6, "เป้าหมาย", "sin(πx)\nและ  x²", TEAL, PALE_TEAL, 14, 20)
    add_card(slide, 9.62, 1.35, 2.85, 1.6, "โมเดล", "Constant\nLinear\nผ่านจุดกำเนิด", ORANGE, PALE_ORANGE, 14, 17)
    add_card(slide, 6.55, 3.25, 2.85, 1.6, "ข้อมูลฝึก", "n = 2\nสุ่ม 50,000 ชุด", GREEN, PALE_TEAL, 14, 18)
    add_card(slide, 9.62, 3.25, 2.85, 1.6, "วิธี fit", "Normal Equation\nnumpy.linalg.lstsq", RED, PALE_ORANGE, 14, 16)
    add_card(slide, 6.55, 5.15, 5.92, 1.0, "Learning curve", "เพิ่ม n = 2 → 100 และเปรียบเทียบ Ein / Eout ที่ σ = 0.0 และ 0.3",
             TEAL, LIGHT, 14, 15)
    add_footer(slide, 2)
    return slide


def add_concept_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "อ่านผลอย่างไร: Bias–Variance Decomposition", 2, "01  บทนำ")
    add_rect(slide, 0.75, 1.45, 11.85, 0.88, LIGHT, radius=False)
    add_text(slide, "E_out  =  Bias²  +  Variance  +  Noise²", 0.95, 1.68, 11.45, 0.38,
             size=25, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    add_text(slide, "นิยามสำคัญ", 0.75, 2.75, 3.5, 0.32, size=18, color=INK, bold=True, margin=0)
    # กรอบสี่เหลี่ยมมุมฉากแบบสีเป็นกลาง ไม่ใช้เส้นกรอบสีจัดจ้าน
    neutral_line = rgb("DEE2E6")
    border_width = Inches(0.06)
    add_rect(slide, 0.75, 3.1, 3.72, 2.28, WHITE, line=neutral_line, line_width=border_width)
    add_rect(slide, 4.78, 3.1, 3.72, 2.28, WHITE, line=neutral_line, line_width=border_width)
    add_rect(slide, 8.81, 3.1, 3.78, 2.28, WHITE, line=neutral_line, line_width=border_width)
    add_rect(slide, 0.75, 3.2, 0.06, 2.1, TEAL)
    add_text(slide, "Bias²  |  ความเอนเอียง", 1.0, 3.22, 3.45, 0.3, size=16, color=TEAL, bold=True, margin=0)
    add_text(slide, "โมเดลเฉลี่ยห่างจาก\nฟังก์ชันเป้าหมายแค่ไหน\n\nสูง = โมเดลเรียบเกินไป", 1.0, 3.62, 3.25, 1.35, size=15, color=INK, margin=0)
    add_rect(slide, 4.78, 3.2, 0.06, 2.1, ORANGE)
    add_text(slide, "Variance  |  ความแปรปรวน", 5.03, 3.22, 3.45, 0.3, size=16, color=ORANGE, bold=True, margin=0)
    add_text(slide, "โมเดลแกว่งมากแค่ไหน\nเมื่อเปลี่ยนชุดข้อมูลฝึก\n\nสูง = ไวต่อข้อมูลมาก", 5.03, 3.62, 3.25, 1.35, size=15, color=INK, margin=0)
    add_rect(slide, 8.81, 3.2, 0.06, 2.1, RED)
    add_text(slide, "Noise²  |  สัญญาณรบกวน", 9.06, 3.22, 3.25, 0.3, size=16, color=RED, bold=True, margin=0)
    add_text(slide, "ความคลาดเคลื่อนที่ลดไม่ได้\nจากข้อมูลที่มี noise\n\nσ เพิ่ม → Eout สูงขึ้น", 9.06, 3.62, 3.1, 1.35, size=15, color=INK, margin=0)
    add_rect(slide, 1.35, 5.72, 10.65, 0.72, GREEN, radius=False)
    add_text(slide, "ประเด็นสำคัญ: เป้าหมายคือ Eout ต่ำ ไม่ใช่แค่ Ein ต่ำบนชุดฝึก",
             1.62, 5.92, 10.1, 0.25, size=16, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 2)
    return slide


def add_results_slide(prs, target, rows, takeaway, page, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, f"ผลลัพธ์เมื่อเป้าหมายคือ {target}", page, kicker)
    add_text(slide, "ค่าเฉลี่ยจากการสุ่มชุดข้อมูลฝึก n = 2 จำนวน 50,000 ชุด",
             0.72, 1.25, 7.5, 0.3, size=13, color=MUTED, margin=0)

    x, y, w, h = 0.62, 1.58, 12.08, 2.78
    table = slide.shapes.add_table(len(rows) + 1, 5, Inches(x), Inches(y), Inches(w), Inches(h)).table
    col_widths = [2.75, 2.05, 2.1, 2.45, 2.73]
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)

    def style_cell(cell, text, fill, color, size=12, bold=False, align=PP_ALIGN.CENTER):
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        for run in p.runs:
            set_font(run, size, color, bold)

    # Single-level header
    headers = ["โมเดล", "Bias² (มือ)", "Variance (มือ)", "Eout (มือ)", "Eout (sim)"]
    for j, val in enumerate(headers):
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        style_cell(table.cell(0, j), val, GREEN, WHITE, 12, True, align)

    min_eout = min(float(row[4]) for row in rows)
    for i, row in enumerate(rows, start=1):
        data_row = i
        is_best = float(row[4]) == min_eout
        base_fill = WHITE if i % 2 else LIGHT
        style_cell(table.cell(data_row, 0), row[0], base_fill, NAVY if is_best else INK, 12, is_best, PP_ALIGN.LEFT)
        for j, value in enumerate(row[1:5], start=1):
            fill = rgb("DFF2E3") if is_best and j == 4 else (rgb("EDF7EF") if j == 4 else base_fill)
            style_cell(table.cell(data_row, j), str(value), fill, NAVY if is_best else INK, 12, is_best)

    add_text(slide, "ข้อสังเกต", 0.62, 4.42, 3.2, 0.3, size=18, color=INK, bold=True, margin=0)
    add_rect(slide, 0.62, 4.82, 0.06, 1.1, RED)
    add_text(slide, takeaway.replace("\n", "  "), 0.9, 4.83, 11.55, 0.86,
             size=16, color=INK, margin=0)
    add_text(slide, "คำนวณมือ: จาก คำนวณมือ.pdf  |  โค้ด: simulation จาก results.json  |  สีเขียวอ่อน = Eout จากโค้ดต่ำสุด",
             0.62, 6.27, 12.0, 0.3, size=10, color=MUTED, margin=0)
    return slide


def add_image_slide(prs, title, image_path, caption, page, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, page, kicker)
    add_contained_image(slide, image_path, 0.62, 1.28, 12.08, 5.35)
    add_text(slide, caption, 0.72, 6.72, 11.9, 0.25, size=10, color=MUTED,
             align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, page)
    return slide


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "สรุป: สิ่งที่ควรจำจากการทดลองนี้", 8, "05  Conclusion")
    add_card(slide, 0.72, 1.42, 3.72, 3.35, "01  โมเดลซับซ้อน", "Linear fit ผ่านจุดฝึก 2 จุดได้ดีมาก\n\nแต่ sin(πx) ทำให้ Variance สูงถึงประมาณ 1.67\n\n→ Ein ต่ำ ไม่ได้แปลว่า Eout ต่ำ", ORANGE, PALE_ORANGE, 15, 19)
    add_card(slide, 4.8, 1.42, 3.72, 3.35, "02  โครงสร้างเป้าหมายสำคัญ", "สำหรับ x² ซึ่งสมมาตรบน [-1, 1]\n\nConstant มี Eout ต่ำสุดประมาณ 0.13\n\n→ เลือกโมเดลให้เข้ากับรูปทรงของข้อมูล", TEAL, PALE_TEAL, 15, 19)
    add_card(slide, 8.88, 1.42, 3.72, 3.35, "03  เพิ่มข้อมูลช่วยลด variance", "เมื่อ n เพิ่มขึ้น Ein และ Eout ค่อย ๆ เข้าใกล้กัน\n\nNoise ทำให้ Eout สูงขึ้น\nแต่ไม่เปลี่ยนแนวโน้มหลัก", GREEN, LIGHT, 15, 19)
    add_rect(slide, 1.2, 5.45, 10.95, 0.95, NAVY, radius=False)
    add_text(slide, "ประโยคสรุปสำหรับการนำเสนอ", 1.5, 5.64, 2.7, 0.22, size=11,
             color=TEAL, bold=True, margin=0)
    add_text(slide, "เราต้องบาลานซ์ความเรียบง่ายของโมเดลกับความสามารถในการจับรูปแบบ\nเพื่อให้โมเดลทำงานได้ดีกับข้อมูลใหม่ ไม่ใช่แค่ข้อมูลที่ใช้ฝึก",
             1.5, 5.91, 10.25, 0.38, size=17, color=WHITE, bold=True, margin=0)
    add_footer(slide, 8)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    add_title_slide(prs)
    add_concept_slide(prs)
    add_results_slide(prs, "sin(πx)", [
        ["Constant", "0.5000", "0.2500", "0.7500", "0.7491"],
        ["Linear", "0.2060", "1.6703", "1.8763", "1.8625"],
        ["Linear ผ่านจุดกำเนิด", "0.2718", "0.2372", "0.5090", "0.5151"],
    ], "Linear มี Bias² ต่ำ\nแต่ Variance สูงมาก\n\nผู้ชนะ: Linear ผ่านจุดกำเนิด\nเพราะ Eout ต่ำสุด", 3, "02  ผลลัพธ์ Bias-Variance")
    add_results_slide(prs, "x²", [
        ["Constant", "0.0889", "0.0444", "0.1333", "0.1346"],
        ["Linear", "0.2000", "0.3333", "0.5333", "0.5414"],
        ["Linear ผ่านจุดกำเนิด", "0.2000", "0.1147", "0.3147", "0.3182"],
    ], "Constant ชนะอย่างชัดเจน\nเพราะ x² สมมาตร\n\nความซับซ้อนเพิ่มขึ้น\nไม่ได้ช่วยให้ Eout ต่ำลง", 4, "02  ผลลัพธ์ Bias-Variance")
    add_image_slide(prs, "ภาพรวม: โมเดลเฉลี่ยและความแกว่ง", PLOTS_DIR / "average_fit.png",
                    "เส้นเขียว = เป้าหมายจริง  |  เส้นประแดง = โมเดลเฉลี่ย  |  แถบแดง = ±1 std  |  เส้นเทา = ตัวอย่างโมเดลจากชุดข้อมูลต่างกัน",
                    5, "03  ภาพประกอบ")
    add_image_slide(prs, "Learning Curve: เมื่อเพิ่มจำนวนข้อมูล", PLOTS_DIR / "learning_curve.png",
                    "เส้นประ = Ein  |  เส้นทึบ = Eout  |  สีน้ำเงิน = σ 0.0  |  สีส้ม = σ 0.3  |  อ่านแนวโน้มจากซ้ายไปขวาเมื่อ n เพิ่มขึ้น",
                    6, "03  ภาพประกอบ")
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build_deck()
