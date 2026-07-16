"""สร้างสไลด์ภาษาไทยสำหรับ Assignment 2: Performance Estimation.

ใช้โทนสีและเลย์เอาต์เดียวกับสไลด์ Week 3 และวางกราฟแบบรักษา aspect ratio
เพื่อไม่ให้ภาพถูกยืดหรือถูกตัดขอบ
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
WEEK4_DIR = BASE_DIR.parent
PLOTS_DIR = WEEK4_DIR / "plots"
OUTPUT = BASE_DIR / "slides_Week4.pptx"

SLIDE_W, SLIDE_H = 13.333, 7.5


def rgb(value):
    return RGBColor.from_string(value.lstrip("#").upper())


NAVY = rgb("2D6A4F")
INK = rgb("212529")
MUTED = rgb("6C757D")
TEAL = rgb("2D6A4F")
ORANGE = rgb("F77F00")
GREEN = rgb("2D6A4F")
RED = rgb("D62828")
LIGHT = rgb("E9ECEF")
PALE_TEAL = rgb("F8F9FA")
PALE_ORANGE = rgb("F8F9FA")
BG = rgb("F8F9FA")
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_font(run, size, color=INK, bold=False, name="Arial"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        set_font(run, size, color, bold)
    return box


def add_rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_footer(slide, page):
    add_rect(slide, 0.62, 6.98, 12.08, 0.012, rgb("DEE2E6"))
    add_text(slide, "Machine Learning — Week 4  |  Performance Estimation", 0.62, 7.08,
             8.5, 0.2, size=8, color=MUTED, margin=0)
    add_text(slide, str(page), 12.12, 7.06, 0.58, 0.22, size=9, color=MUTED,
             bold=True, align=PP_ALIGN.RIGHT, margin=0)


def add_header(slide, title, page, kicker):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_rect(slide, 0, 0, SLIDE_W, 0.42, GREEN)
    add_text(slide, kicker, 0.62, 0.08, 7.0, 0.2, size=10, color=WHITE,
             bold=True, margin=0)
    add_text(slide, title, 0.62, 0.58, 12.0, 0.52, size=25, color=INK,
             bold=True, margin=0)
    add_rect(slide, 0.62, 1.08, 0.86, 0.035, RED)


def add_card(slide, x, y, w, h, title, body, accent=TEAL, fill=LIGHT,
             title_size=14, body_size=16):
    # ใช้กรอบสี่เหลี่ยมมุมฉากแทนกรอบมุมโค้ง เพื่อแบ่งเนื้อหาให้ชัดเจน
    add_rect(slide, x, y, w, h, fill, radius=False)
    add_rect(slide, x, y, 0.08, h, accent)
    add_text(slide, title, x + 0.22, y + 0.16, w - 0.4, 0.28,
             size=title_size, color=accent, bold=True, margin=0)
    add_text(slide, body, x + 0.22, y + 0.52, w - 0.4, h - 0.62,
             size=body_size, color=INK, margin=0)


def add_contained_image(slide, image_path, x, y, w, h):
    """วางภาพให้อยู่ในกรอบโดยคงสัดส่วนเดิมและจัดกึ่งกลาง."""
    with Image.open(image_path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        draw_w = w
        draw_h = w / image_ratio
    else:
        draw_h = h
        draw_w = h * image_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(draw_x), Inches(draw_y),
                             width=Inches(draw_w), height=Inches(draw_h))


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GREEN)
    add_rect(slide, 0.78, 1.52, 1.15, 0.035, WHITE)
    add_text(slide, "MODEL EVALUATION LAB", 0.78, 1.82, 6.0, 0.28,
             size=12, color=WHITE, bold=True, margin=0)
    add_text(slide, "ประเมินประสิทธิภาพโมเดล\nด้วย Resubstitution, Holdout และ K-Fold CV",
             0.78, 2.22, 11.2, 1.25, size=30, color=WHITE, bold=True, margin=0)
    add_text(slide, "ทำไมค่าความผิดพลาดที่วัดได้ จึงไม่เท่ากับ Eout จริงเสมอ",
             0.8, 3.82, 10.8, 0.38, size=18, color=WHITE, margin=0)
    add_text(slide, "Assignment 2  •  Performance Estimation", 0.8, 5.55, 7.0, 0.3,
             size=13, color=WHITE, margin=0)
    add_text(slide, "รายงานผลการทดลอง  |  Machine Learning — Week 4", 0.8, 6.1,
             8.5, 0.25, size=10, color=WHITE, margin=0)
    add_rect(slide, 8.15, 5.05, 0.025, 1.48, WHITE)
    add_text(slide, "ผู้จัดทำ", 8.45, 5.08, 3.8, 0.25, size=11, color=WHITE,
             bold=True, margin=0)
    add_text(slide, "นายสุภมงคล ชอบรัมย์\nนายชินวัตร กิตต๊ะ\nนายกฤตชัย พรายศรี",
             8.45, 5.42, 4.15, 0.9, size=10, color=WHITE, margin=0)
    return slide


def add_setup_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "โจทย์และการตั้งค่าการทดลอง", 2, "01  Setup")
    add_text(slide, "เราจะประมาณ Eout ของโมเดล\nจากข้อมูลชุดเดียวได้อย่างไร?",
             0.72, 1.45, 5.45, 1.1, size=23, color=NAVY, bold=True, margin=0)
    add_text(slide, "เปรียบเทียบวิธีวัด error กับค่า true Eout ที่รู้จากการจำลอง",
             0.72, 2.78, 5.3, 0.45, size=15, color=MUTED, margin=0)
    add_card(slide, 6.55, 1.35, 2.85, 1.6, "เป้าหมาย", "f(x) = sin(πx)\nx ~ U(−1, 1)", TEAL, PALE_TEAL, 14, 18)
    add_card(slide, 9.62, 1.35, 2.85, 1.6, "โมเดล", "Constant\nLinear", ORANGE, PALE_ORANGE, 14, 19)
    add_card(slide, 6.55, 3.25, 2.85, 1.6, "ข้อมูล", "n = 20\nσ = 0.3", GREEN, PALE_TEAL, 14, 19)
    add_card(slide, 9.62, 3.25, 2.85, 1.6, "การทดลอง", "2,000 ชุดข้อมูล\n5-Fold CV", RED, PALE_ORANGE, 14, 18)
    add_card(slide, 6.55, 5.15, 5.92, 1.0, "ค่าอ้างอิง", "true Eout คำนวณบน grid 4,000 จุด และบวก noise² เพื่อสะท้อนข้อมูลใหม่",
             TEAL, LIGHT, 14, 15)
    add_footer(slide, 2)
    return slide


def add_methods_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "สามวิธีประมาณความผิดพลาด", 3, "01  วิธีวัดผล")
    add_text(slide, "แต่ละวิธีใช้ข้อมูลไม่เหมือนกัน จึงมี bias และ variance ต่างกัน",
             0.72, 1.28, 11.8, 0.35, size=15, color=MUTED, margin=0)
    add_card(slide, 0.72, 1.85, 3.72, 2.35, "01  Resubstitution",
             "ฝึกและทดสอบบนข้อมูลชุดเดียวกัน\n\nมักประเมิน error ต่ำเกินจริง\nเพราะโมเดลเคยเห็นข้อมูลแล้ว",
             TEAL, PALE_TEAL, 15, 16)
    add_card(slide, 4.8, 1.85, 3.72, 2.35, "02  Holdout",
             "แบ่งข้อมูลเป็น train / test\n\nสะท้อนข้อมูลใหม่ได้ดีขึ้น\nแต่ผลขึ้นกับสัดส่วนที่แบ่ง",
             ORANGE, PALE_ORANGE, 15, 16)
    add_card(slide, 8.88, 1.85, 3.72, 2.35, "03  K-Fold CV",
             "แบ่งเป็น k folds แล้วสลับ\nชุด validation ให้ครบทุก fold\n\nใช้ข้อมูลคุ้มกว่า holdout",
             GREEN, LIGHT, 15, 16)
    add_rect(slide, 0.72, 4.72, 11.88, 1.08, NAVY)
    add_text(slide, "วัดคุณภาพตัวประมาณ", 1.0, 4.91, 2.55, 0.22, size=12,
             color=WHITE, bold=True, margin=0)
    add_text(slide, "Bias = E[estimate − true Eout]    |    Variance = Var(estimate)    |    MSE = E[(estimate − true Eout)²]",
             1.0, 5.25, 11.1, 0.27, size=15, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 3)
    return slide


def add_plot_slide(prs, title, image_name, caption, page, kicker, note):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, page, kicker)
    add_text(slide, note, 0.72, 1.23, 12.0, 0.3, size=13, color=MUTED,
             align=PP_ALIGN.LEFT, margin=0)
    add_contained_image(slide, PLOTS_DIR / image_name, 0.62, 1.53, 12.08, 5.05)
    add_text(slide, caption, 0.72, 6.68, 11.9, 0.25, size=10, color=MUTED,
             align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, page)
    return slide


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "สรุป: เลือกวิธีวัดผลให้เหมาะกับข้อมูล", 7, "04  Conclusion")
    add_card(slide, 0.72, 1.42, 3.72, 3.35, "01  Resub ต่ำเกินจริง",
             "การทดสอบบนชุดฝึก\nทำให้ error ดูดีเกินไป\n\n→ ใช้ดู fit เบื้องต้น\nแต่อย่าใช้แทน Eout",
             ORANGE, PALE_ORANGE, 15, 18)
    add_card(slide, 4.8, 1.42, 3.72, 3.35, "02  แบ่งข้อมูลมีผล",
             "Holdout ไวต่อสัดส่วน train / test\n\nK-Fold ใช้ข้อมูลคุ้มกว่า\nและมักมี variance ต่ำลงเมื่อ k เพิ่ม",
             TEAL, PALE_TEAL, 15, 18)
    add_card(slide, 8.88, 1.42, 3.72, 3.35, "03  ข้อมูลช่วยลด variance",
             "เพิ่ม n ทำให้ค่าประมาณนิ่งขึ้น\n\nnoise สูงขึ้น → variance สูงขึ้น\nโดยเฉพาะ Linear + Holdout",
             GREEN, LIGHT, 15, 18)
    add_rect(slide, 1.2, 5.45, 10.95, 0.95, NAVY)
    add_text(slide, "ประโยคสรุปสำหรับการนำเสนอ", 1.5, 5.64, 2.7, 0.22,
             size=11, color=TEAL, bold=True, margin=0)
    add_text(slide, "การวัดผลที่ดีต้องเลียนแบบข้อมูลใหม่ให้มากที่สุด และต้องพิจารณา bias กับ variance ไปพร้อมกัน",
             1.5, 5.91, 10.25, 0.38, size=17, color=WHITE, bold=True, margin=0)
    add_footer(slide, 7)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    add_title_slide(prs)
    add_setup_slide(prs)
    add_methods_slide(prs)
    add_plot_slide(
        prs,
        "Bias / Variance / MSE ของตัวประมาณ",
        "part2.png",
        "แกนตั้ง = estimate − true Eout  |  เส้นดำ = 0  |  สามเหลี่ยมสีเขียว = ค่าเฉลี่ย",
        4,
        "02  ผลการทดลอง",
        "Resubstitution มีแนวโน้มติดลบ เพราะประเมิน error ต่ำกว่าความจริง",
    )
    add_plot_slide(
        prs,
        "ผลของสัดส่วน Holdout และจำนวน fold",
        "part3.png",
        "ซ้าย: เปลี่ยนสัดส่วน train ของ Holdout  |  ขวา: เปลี่ยน k ของ K-Fold  |  ค่าที่ต่ำกว่ามักหมายถึง variance น้อยกว่า",
        5,
        "03  Sensitivity Analysis",
        "Holdout ที่ใช้ข้อมูลฝึกน้อยเกินไปมี variance สูงมาก ส่วน K-Fold นิ่งขึ้นเมื่อ k เพิ่ม",
    )
    add_plot_slide(
        prs,
        "ผลของจำนวนข้อมูลและระดับ Noise",
        "part4.png",
        "แถวบน = Constant  |  แถวล่าง = Linear  |  น้ำเงิน = Resub  |  ส้ม = Holdout  |  เขียว = K-Fold",
        6,
        "03  Sensitivity Analysis",
        "เมื่อ n เพิ่ม variance ลดลง; เมื่อ σ เพิ่ม variance สูงขึ้น โดย Holdout และ Linear ได้รับผลชัดที่สุด",
    )
    add_summary_slide(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build_deck()
